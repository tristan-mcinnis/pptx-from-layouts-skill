#!/usr/bin/env python3
"""
catalog.py — Step 1 of the PPTX-from-layouts pipeline: PROFILE.

Given any .pptx template, this emits two artifacts that the next two steps
consume:

  1. <name>-config.json     — the machine-readable layout map (Step 3 / render
                              reads this via `generate.py --config`). Produced by
                              the sibling `pptx-from-layouts` profiler so the
                              format stays identical to what the engine expects.
  2. <name>-layout-catalog.md — the human-readable "menu" of layout names. This
                              is what an author reads in Step 2 to know which
                              `[HINT: layout_name]` values are valid for THIS
                              template, and what each layout is for.

Usage:
    python catalog.py /path/to/template.pptx --name my-template --output-dir ./build
    python catalog.py /path/to/template.pptx                 # name/dir inferred

Dependencies: python-pptx (already required by the render engine). The config
step shells out to the sibling pptx-from-layouts skill, which must ship
alongside this skill in the same repo (../pptx-from-layouts).
"""

import argparse
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx",
          file=sys.stderr)
    sys.exit(1)

_SCRIPT_DIR = Path(__file__).resolve().parent
_SKILL_DIR = _SCRIPT_DIR.parent
_SKILLS_ROOT = _SKILL_DIR.parent  # .claude/skills/
# The render engine ships beside us in the same repo.
_RENDER_PROFILE = _SKILLS_ROOT / "pptx-from-layouts" / "scripts" / "profile.py"

# python-pptx placeholder type → friendly label
_PH_KIND = {
    0: "title", 1: "body", 2: "center-title", 3: "subtitle", 4: "vertical-body",
    7: "object", 8: "chart", 9: "table", 10: "clip-art", 11: "diagram",
    12: "media", 13: "picture", 14: "header", 15: "footer", 16: "date",
    17: "slide-number", 18: "vertical-object", 19: "vertical-title",
}


def _ph_label(ph) -> str:
    try:
        return _PH_KIND.get(int(ph.placeholder_format.type), str(ph.placeholder_format.type))
    except Exception:
        return "?"


def _suggest_use(name: str, kinds: list[str]) -> str:
    """Heuristic 'when to use' line from the layout name + placeholder makeup."""
    n = name.lower()
    has_pic = "picture" in kinds
    if "master" in n or "base" in n:
        return "Branding / logo slide. Usually slide 1 — kept from the template."
    if "title" in n and ("cover" in n or "centered" not in n) and "section" not in n:
        return "Cover / title slide (project, client, date)."
    if "section" in n or "divider" in n:
        return "Section break — single title only. Use sparingly."
    if "contact" in n:
        return "Closing / contact slide."
    if "column" in n or "grid" in n:
        digits = "".join(c for c in n if c.isdigit())
        cols = digits[:1] or "N"
        return f"{cols}-up parallel items (cards / comparison / framework)" + (" with images." if has_pic else ".")
    if "quote" in n:
        return "Pull-quote / hero quote."
    if "image" in n and has_pic:
        return "Body content with a supporting image."
    if has_pic:
        return "Content with an image placeholder."
    return "General content slide (title + body)."


def _placeholders_by_index(prs) -> dict[int, list[str]]:
    """Map each layout index → list of friendly placeholder labels."""
    out: dict[int, list[str]] = {}
    for idx, layout in enumerate(prs.slide_layouts):
        out[idx] = [_ph_label(ph) for ph in layout.placeholders]
    return out


def build_catalog(pptx_path: Path, template_name: str, config_path: Path | None = None) -> str:
    import json as _json

    prs = Presentation(str(pptx_path))
    emu_w, emu_h = prs.slide_width, prs.slide_height
    try:
        dims = f"{Emu(emu_w).inches:.2f}\" x {Emu(emu_h).inches:.2f}\""
    except Exception:
        dims = "unknown"
    master_name = (prs.slide_masters[0].name or "(unnamed)") if prs.slide_masters else "(unnamed)"
    ph_by_idx = _placeholders_by_index(prs)

    # Load the curated capability map (if config was generated). Every layout in
    # here is one the render engine knows how to populate — so these become the
    # first-class, guaranteed-renderable hints.
    curated: list[tuple[int, str, str]] = []  # (index, layout_name, use_case)
    if config_path and config_path.exists():
        try:
            cfg = _json.loads(config_path.read_text(encoding="utf-8"))
            seen = set()
            for cap_key, cap in cfg.get("layout_mappings", {}).items():
                lname = cap.get("layout_name", "")
                lidx = cap.get("layout_index", -1)
                if lname and lname not in seen:
                    seen.add(lname)
                    curated.append((lidx, lname, cap.get("use_case", "")))
            curated.sort(key=lambda t: t[0])
        except Exception:
            curated = []

    lines: list[str] = []
    lines.append(f"# Layout catalog — {template_name}")
    lines.append("")
    lines.append("> Generated by **Step 1 (pptx-profile)**. This is the menu of layout")
    lines.append("> names you can target in **Step 2** with `[HINT: layout_name]`, and that")
    lines.append("> **Step 3** renders into. Copy a name from the **Hint** column verbatim.")
    lines.append("")
    lines.append(f"- **Template:** `{pptx_path.name}`")
    lines.append(f"- **Slide size:** {dims}")
    lines.append(f"- **Slide master:** {master_name}")
    lines.append(f"- **Layouts in file:** {len(prs.slide_layouts)}"
                 + (f"  ·  **Hintable (mapped):** {len(curated)}" if curated else ""))
    lines.append("")

    if curated:
        lines.append("## Hintable layouts (use these)")
        lines.append("")
        lines.append("Every name here is guaranteed renderable — the engine knows how to")
        lines.append("fill it. These are your safe `[HINT:]` values.")
        lines.append("")
        lines.append("| Idx | Hint (use this name) | Placeholders | When to use |")
        lines.append("|----:|----------------------|--------------|-------------|")
        for lidx, lname, use in curated:
            hint = lname.strip().lower().replace(" ", "-")
            kinds = ph_by_idx.get(lidx, [])
            ph_summary = ", ".join(kinds) if kinds else "— (none)"
            use = use or _suggest_use(lname, kinds)
            lines.append(f"| {lidx} | `{hint}` | {ph_summary} | {use} |")
        lines.append("")
        lines.append("<details><summary>All %d layouts in the file (reference — may need "
                     "config tuning before hinting)</summary>" % len(prs.slide_layouts))
        lines.append("")
        header = "| Idx | Layout name | Placeholders |"
        sep = "|----:|-------------|--------------|"
    else:
        lines.append("> No render config was found, so every raw layout is listed. Hints")
        lines.append("> resolve best once `*-config.json` exists (re-run without `--catalog-only`).")
        lines.append("")
        header = "| Idx | Hint (use this name) | Placeholders | When to use |"
        sep = "|----:|----------------------|--------------|-------------|"

    lines.append(header)
    lines.append(sep)
    for idx, layout in enumerate(prs.slide_layouts):
        name = layout.name or f"layout-{idx}"
        hint = name.strip().lower().replace(" ", "-")
        kinds = ph_by_idx.get(idx, [])
        ph_summary = ", ".join(kinds) if kinds else "— (none)"
        if curated:
            lines.append(f"| {idx} | `{hint}` | {ph_summary} |")
        else:
            lines.append(f"| {idx} | `{hint}` | {ph_summary} | {_suggest_use(name, kinds)} |")
    if curated:
        lines.append("")
        lines.append("</details>")

    lines.append("")
    lines.append("## How to use this in Step 2")
    lines.append("")
    lines.append("```markdown")
    lines.append("# Slide 1: <a claim, not a label>")
    lines.append("")
    lines.append("[HINT: <copy a Hint value from the table above>]")
    lines.append("")
    lines.append("- Your bullet / column / table content")
    lines.append("```")
    lines.append("")
    lines.append("Layouts with a **picture** placeholder accept `[Image: path]`. Layouts")
    lines.append("named `column-*` / `grid-*` accept `[Column N: Header]` blocks. If you")
    lines.append("omit `[HINT:]`, the engine auto-selects a layout from your `**Visual:**`")
    lines.append("type — hints just make the choice explicit and template-accurate.")
    lines.append("")
    return "\n".join(lines)


def generate_config(pptx_path: Path, name: str, out_dir: Path) -> Path | None:
    """Shell out to the sibling render engine's profiler to emit <name>-config.json."""
    if not _RENDER_PROFILE.exists():
        print(f"WARNING: render engine profiler not found at {_RENDER_PROFILE}.\n"
              f"         Skipping config.json — Step 3 needs it for custom templates.",
              file=sys.stderr)
        return None
    cmd = [
        sys.executable, str(_RENDER_PROFILE), str(pptx_path.resolve()),
        "--name", name, "--generate-config", "--output-dir", str(out_dir.resolve()),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"WARNING: config generation failed:\n{result.stderr or result.stdout}",
              file=sys.stderr)
        return None
    config_path = out_dir / f"{name}-config.json"
    return config_path if config_path.exists() else None


def main():
    ap = argparse.ArgumentParser(description="Step 1: profile a .pptx into a layout catalog + config.")
    ap.add_argument("template", help="Path to the .pptx template")
    ap.add_argument("--name", "-n", help="Template name (default: from filename)")
    ap.add_argument("--output-dir", "-o", default=".", help="Where to write artifacts")
    ap.add_argument("--catalog-only", action="store_true",
                    help="Emit only the markdown catalog (skip config.json)")
    args = ap.parse_args()

    pptx_path = Path(args.template)
    if not pptx_path.exists():
        print(f"ERROR: template not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    name = args.name or pptx_path.stem
    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    # 1. Machine config for the render step (also drives the curated catalog).
    config_path = None
    if not args.catalog_only:
        config_path = generate_config(pptx_path, name, out_dir)

    # 2. Human-readable catalog (config-driven when available).
    catalog_md = build_catalog(pptx_path, name, config_path)
    catalog_path = out_dir / f"{name}-layout-catalog.md"
    catalog_path.write_text(catalog_md, encoding="utf-8")
    print(f"✓ Layout catalog: {catalog_path}")
    if config_path:
        print(f"✓ Render config:  {config_path}")

    print("\nNext steps:")
    print(f"  Step 2 — author slides with [HINT:] using {catalog_path.name}")
    print(f"           lint:  python ../pptx-author/scripts/lint_hints.py slides.md "
          f"--config {config_path.name if config_path else '<config>.json'}")
    print(f"  Step 3 — render: python ../pptx-from-layouts/scripts/generate.py slides.md "
          f"-o deck.pptx --template {pptx_path.name} "
          f"--config {config_path.name if config_path else '<config>.json'} --validate")


if __name__ == "__main__":
    main()
